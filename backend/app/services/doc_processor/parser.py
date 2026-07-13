# backend/app/services/doc_processor/parser.py
"""
Document parser service using unstructured.io element types.

Note: unstructured.partition.* functions rely on libmagic (python-magic) which
causes a segfault on this Windows environment. Instead, we implement format-specific
parsing directly and wrap results in unstructured element objects so downstream
code can work with the standard element API.
"""
import re
from pathlib import Path
from typing import Literal

from unstructured.documents.elements import (
    ElementMetadata,
    Image,
    NarrativeText,
    Table,
    Text,
    Title,
)

ElementType = Literal["text", "table", "image", "title"]

_IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif"}

# Excel 每个 chunk 最多包含的数据行数
_XLSX_ROWS_PER_CHUNK = 20

# Patterns used to classify text blocks
_TITLE_PATTERN = re.compile(
    r"^(第[一二三四五六七八九十百]+[章节部分]|"
    r"Chapter\s+\d+|"
    r"Section\s+\d+|"
    r"\d+\.\s+\S|"
    r"[A-Z][A-Z\s]{2,50}$)",
    re.UNICODE,
)


def parse_document(file_path: str) -> list[dict]:
    """
    解析文档，返回结构化元素列表。
    每个元素：{
        "content": str,         # 文本内容（图片为空字符串，等待 OCR）
        "element_type": str,    # text / table / image / title
        "page_number": int,     # 页码（从1开始）
        "raw_element": object,  # 原始 unstructured 元素（供后续处理）
    }
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    suffix = path.suffix.lower()

    if suffix == ".txt":
        return _parse_txt(path)
    elif suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix in (".docx", ".doc"):
        return _parse_docx(path)
    elif suffix in (".xlsx", ".xls"):
        return _parse_xlsx(path)
    elif suffix == ".pptx":
        return _parse_pptx(path)
    elif suffix in _IMAGE_SUFFIXES:
        return _parse_image_file(path)
    else:
        # Generic fallback: try to read as text
        return _parse_txt(path)


# ---------------------------------------------------------------------------
# Format-specific parsers
# ---------------------------------------------------------------------------

def _parse_txt(path: Path) -> list[dict]:
    """Parse plain text files by splitting on blank lines."""
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = path.read_text(encoding="gbk", errors="replace")

    paragraphs = re.split(r"\n{2,}", content.strip())
    result = []
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        el_type = _classify_text(para)
        el = _make_element(el_type, para, page_number=1)
        result.append(_to_dict(el, el_type, page_number=1))
    return result


def _parse_pdf(path: Path) -> list[dict]:
    """Parse PDF using pypdf (pure Python, no libmagic dependency).
    Only extracts text layers; image-only (scanned) PDFs will return empty list.
    """
    try:
        import pypdf  # type: ignore
    except ImportError:
        try:
            import PyPDF2 as pypdf  # type: ignore
        except ImportError:
            raise ImportError(
                "pypdf or PyPDF2 is required for PDF parsing. "
                "Install with: pip install pypdf"
            )

    result = []
    reader = pypdf.PdfReader(str(path))
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for para in re.split(r"\n{2,}", text.strip()):
            para = para.strip()
            if not para:
                continue
            el_type = _classify_text(para)
            el = _make_element(el_type, para, page_number=page_num)
            result.append(_to_dict(el, el_type, page_number=page_num))
    return result


def _convert_doc_to_docx(path: Path) -> Path:
    """Convert legacy .doc to .docx using LibreOffice. Returns path to .docx file."""
    import subprocess
    import shutil
    if not shutil.which("libreoffice") and not shutil.which("soffice"):
        raise RuntimeError(
            "不支持 .doc 格式，请将文件转换为 .docx 后重新上传，"
            "或在服务器上安装 LibreOffice（apt-get install libreoffice）"
        )
    cmd = shutil.which("libreoffice") or shutil.which("soffice")
    out_dir = path.parent
    subprocess.run(
        [cmd, "--headless", "--convert-to", "docx", "--outdir", str(out_dir), str(path)],
        check=True,
        capture_output=True,
        timeout=60,
    )
    docx_path = path.with_suffix(".docx")
    if not docx_path.exists():
        raise RuntimeError(f"LibreOffice 转换失败，未生成 .docx 文件")
    return docx_path


def _parse_docx(path: Path) -> list[dict]:
    """Parse DOCX using python-docx. Converts .doc via LibreOffice if needed."""
    try:
        import docx  # type: ignore
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX parsing. "
            "Install with: pip install python-docx"
        )

    # .doc (旧格式) 需要先转换
    if path.suffix.lower() == ".doc":
        path = _convert_doc_to_docx(path)

    doc = docx.Document(str(path))
    result = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading" in style or "title" in style:
            el_type: ElementType = "title"
        else:
            el_type = _classify_text(text)
        el = _make_element(el_type, text, page_number=1)
        result.append(_to_dict(el, el_type, page_number=1))

    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_text = "\n".join(["\t".join(row) for row in rows])
        el = Table(text=table_text)
        el.metadata.page_number = 1
        result.append(_to_dict(el, "table", page_number=1))

    return result


def _parse_xlsx(path: Path) -> list[dict]:
    """Parse XLSX using openpyxl.

    每个 sheet 按行分块（每块最多 _XLSX_ROWS_PER_CHUNK 行），
    每块开头携带 sheet 名称和列名作为上下文前缀，避免单个巨型 chunk 稀释向量语义。
    """
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise ImportError(
            "openpyxl is required for XLSX parsing. "
            "Install with: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    result = []

    for sheet_idx, sheet in enumerate(wb.worksheets):
        all_rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            all_rows.append([str(cell) if cell is not None else "" for cell in row])

        if not all_rows:
            continue

        # 第一行视为列名
        header_row = all_rows[0]
        data_rows = all_rows[1:]
        header_str = "\t".join(header_row)
        # 列名前缀：[工作表: Sheet1] 列名: A | B | C
        col_prefix = f"[工作表: {sheet.title}] 列名: {' | '.join(h for h in header_row if h)}"

        if not data_rows:
            # 只有列名行
            el = Table(text=col_prefix)
            el.metadata.page_number = sheet_idx + 1
            result.append(_to_dict(el, "table", page_number=sheet_idx + 1))
            continue

        # 按 _XLSX_ROWS_PER_CHUNK 行拆分，每块带 col_prefix + header_str
        for chunk_start in range(0, len(data_rows), _XLSX_ROWS_PER_CHUNK):
            chunk_rows = data_rows[chunk_start: chunk_start + _XLSX_ROWS_PER_CHUNK]
            rows_text = "\n".join("\t".join(row) for row in chunk_rows)
            table_text = f"{col_prefix}\n{header_str}\n{rows_text}"
            el = Table(text=table_text)
            el.metadata.page_number = sheet_idx + 1
            result.append(_to_dict(el, "table", page_number=sheet_idx + 1))

    return result


def _parse_pptx(path: Path) -> list[dict]:
    """Parse PPTX using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX parsing. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation(str(path))
    result = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                el_type = _classify_text(text)
                el = _make_element(el_type, text, page_number=slide_num)
                result.append(_to_dict(el, el_type, page_number=slide_num))
    return result


def _parse_image_file(path: Path) -> list[dict]:
    """解析独立图片文件（.jpg/.png 等）。

    优先调用 doubao-seed-1.6 视觉模型生成图片描述；
    未配置 vision_api_key 时 fallback 到 EasyOCR 提取图中文字。
    """
    from app.services.doc_processor.ocr import analyze_image_with_vision, extract_text_from_image

    image_bytes = path.read_bytes()

    # 1. 尝试 vision 模型（doubao-seed-1.6）
    text = analyze_image_with_vision(image_bytes, filename=path.name)

    # 2. fallback：EasyOCR 识别图中文字
    if not text:
        text = extract_text_from_image(image_bytes)

    # 3. 兜底提示
    if not text:
        text = f"[图片文件: {path.name}，未能提取内容]"

    el = NarrativeText(text=text)
    el.metadata.page_number = 1
    return [_to_dict(el, "text", page_number=1)]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _classify_text(text: str) -> ElementType:
    """Heuristic classification of a text block."""
    stripped = text.strip()
    if not stripped:
        return "text"
    # Short lines that look like headings
    if len(stripped) <= 80 and _TITLE_PATTERN.match(stripped):
        return "title"
    return "text"


def _make_element(el_type: ElementType, text: str, page_number: int):
    """Create an unstructured element of the given type."""
    if el_type == "title":
        el = Title(text=text)
    elif el_type == "table":
        el = Table(text=text)
    elif el_type == "image":
        el = Image(text=text)
    else:
        el = NarrativeText(text=text)
    el.metadata.page_number = page_number
    return el


def _to_dict(el, el_type: ElementType, page_number: int) -> dict:
    content = ""
    if el_type != "image":
        content = str(el.text) if el.text else ""

    page = page_number
    if hasattr(el, "metadata") and el.metadata:
        page = getattr(el.metadata, "page_number", page_number) or page_number

    return {
        "content": content,
        "element_type": el_type,
        "page_number": page,
        "raw_element": el,
    }
